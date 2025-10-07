import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import io
import camelot
import torch
from transformers import CLIPProcessor, CLIPModel
from sentence_transformers import SentenceTransformer
from typing import List, Tuple

# Initialize models
text_model = SentenceTransformer('all-MiniLM-L6-v2')
clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")

def extract_pdf(file_path: str) -> Tuple[List[str], List[bytes], List[str]]:
    doc = fitz.open(file_path)
    text_blocks, images, tables = [], [], []

    # Extract text
    for page in doc:
        text = page.get_text()
        if text.strip():
            text_blocks.append(text)

    # Extract images
    for i in range(len(doc)):
        for img_index, img in enumerate(doc[i].get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            images.append(image_bytes)

    # Extract tables with Camelot
    try:
        tables_found = camelot.read_pdf(file_path, pages='all')
        for table in tables_found:
            tables.append(table.df.to_csv(index=False))
    except:
        pass

    return text_blocks, images, tables

def extract_ppt(file_path: str) -> Tuple[List[str], List[bytes], List[str]]:
    prs = Presentation(file_path)
    text_blocks, images, tables = [], [], []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    text_blocks.append(text)
            if shape.shape_type == 19:  # Table
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append('\n'.join([','.join(r) for r in table_data]))
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                images.append(image_bytes)

    return text_blocks, images, tables

def vectorize_texts(texts: List[str]) -> List[List[float]]:
    if not texts:
        return []
    return text_model.encode(texts, convert_to_tensor=True).tolist()

def vectorize_images(image_bytes_list: List[bytes]) -> List[List[float]]:
    vectors = []
    for img_bytes in image_bytes_list:
        image = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        inputs = clip_processor(images=image, return_tensors="pt")
        with torch.no_grad():
            vec = clip_model.get_image_features(**inputs)
            vec = vec / vec.norm(dim=-1, keepdim=True)
            vectors.append(vec[0].tolist())
    return vectors
